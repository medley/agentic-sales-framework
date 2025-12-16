#!/usr/bin/env python3
"""
Extract text inventory from PowerPoint presentations.

Scans slides for shapes containing text and returns their positions
for use in thumbnail generation with placeholder outlines.
"""

from dataclasses import dataclass
from pathlib import Path
from typing import Dict

from pptx import Presentation
from pptx.util import Inches


@dataclass
class ShapeData:
    """Data class representing a shape's position and dimensions in inches."""
    left: float
    top: float
    width: float
    height: float
    text: str = ""


def extract_text_inventory(
    pptx_path: Path, prs: Presentation = None
) -> Dict[str, Dict[str, ShapeData]]:
    """
    Extract all text-containing shapes from a PowerPoint presentation.

    Args:
        pptx_path: Path to the PowerPoint file
        prs: Optional pre-loaded Presentation object

    Returns:
        Dictionary mapping slide keys ("slide-1", "slide-2", etc.) to
        dictionaries of shape data. Each shape dict maps shape identifiers
        to ShapeData objects with position/dimension info.

    Example return:
        {
            "slide-1": {
                "shape-0": ShapeData(left=0.5, top=1.0, width=9.0, height=1.5, text="Title"),
                "shape-1": ShapeData(left=0.5, top=3.0, width=9.0, height=4.0, text="Content")
            },
            "slide-2": {...}
        }
    """
    if prs is None:
        prs = Presentation(str(pptx_path))

    inventory = {}
    EMU_PER_INCH = 914400.0

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_key = f"slide-{slide_idx}"
        shapes_dict = {}
        shape_counter = 0

        for shape in slide.shapes:
            # Check if shape has text
            if not hasattr(shape, "text_frame"):
                continue

            try:
                text = shape.text_frame.text.strip()
            except Exception:
                continue

            # Only include shapes that have actual text content
            if not text:
                continue

            # Convert EMU (English Metric Units) to inches
            left = (shape.left or 0) / EMU_PER_INCH
            top = (shape.top or 0) / EMU_PER_INCH
            width = (shape.width or 0) / EMU_PER_INCH
            height = (shape.height or 0) / EMU_PER_INCH

            shape_key = f"shape-{shape_counter}"
            shapes_dict[shape_key] = ShapeData(
                left=left,
                top=top,
                width=width,
                height=height,
                text=text[:100]  # Truncate for memory efficiency
            )
            shape_counter += 1

        if shapes_dict:
            inventory[slide_key] = shapes_dict

    return inventory


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Usage: python inventory.py <presentation.pptx>")
        sys.exit(1)

    pptx_path = Path(sys.argv[1])
    if not pptx_path.exists():
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)

    inventory = extract_text_inventory(pptx_path)

    print(f"Found text in {len(inventory)} slides:")
    for slide_key, shapes in inventory.items():
        print(f"\n{slide_key}:")
        for shape_key, data in shapes.items():
            print(f"  {shape_key}: '{data.text[:50]}...' at ({data.left:.1f}, {data.top:.1f})")
